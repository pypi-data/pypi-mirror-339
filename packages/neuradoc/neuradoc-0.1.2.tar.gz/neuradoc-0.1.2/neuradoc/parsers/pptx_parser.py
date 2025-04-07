"""
Parser for Microsoft PowerPoint presentations.
"""

import logging
from pathlib import Path
import io

from pptx import Presentation
from PIL import Image

from neuradoc.parsers import BaseParser
from neuradoc.models.element import Element, ElementType

logger = logging.getLogger(__name__)


class Parser(BaseParser):
    """Parser for Microsoft PowerPoint presentations."""
    
    def parse(self, file_path):
        """
        Parse a PowerPoint presentation.
        
        Args:
            file_path (str): Path to the PPTX file
            
        Returns:
            dict: Parsed content with metadata
        """
        self.validate_file(file_path)
        
        try:
            # Open the presentation
            prs = Presentation(file_path)
            
            # Extract document metadata
            metadata = {
                'title': Path(file_path).stem,
                'slide_count': len(prs.slides),
                'file_path': file_path,
                'file_type': 'pptx'
            }
            
            # Extract elements
            elements = []
            
            # Process each slide
            for slide_index, slide in enumerate(prs.slides):
                slide_elements = self._process_slide(slide, slide_index)
                elements.extend(slide_elements)
            
            return {
                'metadata': metadata,
                'elements': elements
            }
            
        except Exception as e:
            logger.error(f"Error parsing PPTX {file_path}: {e}")
            raise
    
    def _process_slide(self, slide, slide_index):
        """
        Process a single slide and extract elements.
        
        Args:
            slide: PowerPoint slide object
            slide_index: Index of the slide
            
        Returns:
            list: Extracted elements from this slide
        """
        elements = []
        
        # Extract slide title if available
        if slide.shapes.title and slide.shapes.title.text:
            elements.append(Element(
                element_type=ElementType.TEXT,
                content=slide.shapes.title.text,
                metadata={'type': 'title'},
                position={'slide': slide_index, 'type': 'title'}
            ))
        
        # Process each shape in the slide
        for shape_index, shape in enumerate(slide.shapes):
            # Extract text from text boxes
            if shape.has_text_frame:
                text = ""
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text
                
                if text.strip():
                    elements.append(Element(
                        element_type=ElementType.TEXT,
                        content=text,
                        position={'slide': slide_index, 'shape': shape_index}
                    ))
            
            # Extract tables
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = []
                    for cell in row.cells:
                        cell_text = ""
                        for paragraph in cell.text_frame.paragraphs:
                            cell_text += paragraph.text + "\n"
                        row_data.append(cell_text.strip())
                    table_data.append(row_data)
                
                if table_data:
                    elements.append(Element(
                        element_type=ElementType.TABLE,
                        content=table_data,
                        position={'slide': slide_index, 'shape': shape_index}
                    ))
            
            # Extract images
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    # Extract image data
                    if hasattr(shape, 'image') and shape.image:
                        image_data = shape.image.blob
                        image = Image.open(io.BytesIO(image_data))
                        
                        elements.append(Element(
                            element_type=ElementType.IMAGE,
                            content=image,
                            position={'slide': slide_index, 'shape': shape_index}
                        ))
                except Exception as e:
                    logger.warning(f"Could not extract image: {e}")
            
            # Check for charts (diagrams)
            if shape.shape_type == 3:  # MSO_SHAPE_TYPE.CHART
                elements.append(Element(
                    element_type=ElementType.DIAGRAM,
                    content=f"Chart in slide {slide_index + 1}",
                    position={'slide': slide_index, 'shape': shape_index}
                ))
        
        return elements
