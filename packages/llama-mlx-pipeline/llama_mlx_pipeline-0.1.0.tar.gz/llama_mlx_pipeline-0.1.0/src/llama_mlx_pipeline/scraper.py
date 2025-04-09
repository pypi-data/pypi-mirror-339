import base64
import csv
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import time
import uuid
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union
from urllib.parse import urlparse

import docx

# Document processing
import fitz  # PyMuPDF
import frontmatter
import markdown

# MLX integration
import mlx.core as mx
import numpy as np
import pptx
import pytesseract

# Core libraries
import requests
from bs4 import BeautifulSoup
from PIL import Image

# For Whisper integration
# We'll implement a minimal wrapper for MLX-based Whisper


class WhisperMLX:
    """
    A simplified wrapper for MLX-based Whisper model.
    This is a placeholder - actual implementation would need to match the mlx-examples/whisper implementation.
    """

    def __init__(self, model_path: str):
        """
        Initialize the Whisper model.

        Args:
            model_path: Path to the MLX Whisper model
        """
        self.model_path = model_path
        self.model = None
        self.processor = None
        self._load_model()

    def _load_model(self):
        """Load the MLX Whisper model."""
        try:
            # This is a placeholder for actual model loading
            # In a real implementation, this would use MLX to load the model
            print(f"Loading MLX Whisper model from {self.model_path}")

            # Check if we have the mlx_whisper module
            try:
                import mlx_whisper

                self.model, self.processor = mlx_whisper.load_model(self.model_path)
                print("Loaded Whisper model with mlx_whisper")
            except ImportError:
                # Fallback to a simple implementation
                print("mlx_whisper not found, using simplified implementation")
                self._load_custom_model()

        except Exception as e:
            print(f"Failed to load Whisper model: {str(e)}")
            raise

    def _load_custom_model(self):
        """
        Custom model loading for when mlx_whisper is not available.
        This is a placeholder implementation.
        """
        # In a real implementation, this would load the model architecture and weights
        pass

    def transcribe(self, audio_path: str, **kwargs) -> Dict[str, Any]:
        """
        Transcribe audio file to text.

        Args:
            audio_path: Path to the audio file
            **kwargs: Additional arguments

        Returns:
            Dict: Transcription result with text and segments
        """
        print(f"Transcribing {audio_path}")

        # This is a placeholder for actual transcription
        # In a real implementation, this would:
        # 1. Load and preprocess the audio
        # 2. Run the MLX Whisper model
        # 3. Post-process the results

        # For now, we'll try to use ffmpeg to convert the audio to the right format
        # and then simulate the Whisper output format

        # Check if we need to convert the audio
        temp_file = None
        if not audio_path.endswith(".wav"):
            print("Converting audio to WAV format")
            temp_file = tempfile.NamedTemporaryFile(suffix=".wav", delete=False)
            temp_file.close()

            try:
                subprocess.run(
                    [
                        "ffmpeg",
                        "-i",
                        audio_path,
                        "-ar",
                        "16000",
                        "-ac",
                        "1",
                        "-c:a",
                        "pcm_s16le",
                        temp_file.name,
                    ],
                    check=True,
                    stdout=subprocess.PIPE,
                    stderr=subprocess.PIPE,
                )
                audio_path = temp_file.name
            except Exception as e:
                print(f"Failed to convert audio: {str(e)}")
                os.unlink(temp_file.name)
                raise

        # This is where we'd actually run the model
        # For now, return a placeholder result

        result = {
            "text": "[This is a placeholder transcription - MLX Whisper model would generate actual content]",
            "segments": [
                {
                    "id": 0,
                    "start": 0.0,
                    "end": 10.0,
                    "text": "[This is a placeholder segment - MLX Whisper model would generate actual content]",
                }
            ],
            "language": "en",
        }

        # Clean up temp file if created
        if temp_file and os.path.exists(temp_file.name):
            os.unlink(temp_file.name)

        return result


def scrape_text(source: str) -> str:
    """
    Scrape text from a text file.

    Args:
        source: Path to text file

    Returns:
        str: Text content
    """
    with open(source, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def scrape_html(source: str) -> str:
    """
    Scrape text from an HTML file.

    Args:
        source: Path to HTML file

    Returns:
        str: Extracted text content
    """
    with open(source, "r", encoding="utf-8", errors="replace") as f:
        html = f.read()

    soup = BeautifulSoup(html, "html.parser")

    # Remove script and style elements
    for script in soup(["script", "style"]):
        script.extract()

    # Extract text
    text = soup.get_text()

    # Remove extra whitespace
    lines = (line.strip() for line in text.splitlines())
    chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
    text = "\n".join(chunk for chunk in chunks if chunk)

    return text


def scrape_pdf(source: str) -> str:
    """
    Scrape text from a PDF file.

    Args:
        source: Path to PDF file

    Returns:
        str: Extracted text content
    """
    text = ""

    try:
        doc = fitz.open(source)

        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text += page.get_text() + "\n\n"

        doc.close()
    except Exception as e:
        print(f"Error extracting text from PDF: {str(e)}")
        # Try OCR as a fallback
        text = ocr_pdf(source)

    return text


def ocr_pdf(source: str) -> str:
    """
    OCR a PDF file using pytesseract.

    Args:
        source: Path to PDF file

    Returns:
        str: Extracted text content
    """
    text = ""

    try:
        doc = fitz.open(source)

        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            pix = page.get_pixmap()

            # Save the image to a temporary file
            img_path = f"{tempfile.gettempdir()}/page_{page_num}.png"
            pix.save(img_path)

            # OCR the image
            image = Image.open(img_path)
            page_text = pytesseract.image_to_string(image)
            text += page_text + "\n\n"

            # Clean up
            image.close()
            os.remove(img_path)

        doc.close()
    except Exception as e:
        print(f"Error during OCR: {str(e)}")

    return text


def scrape_docx(source: str) -> str:
    """
    Scrape text from a DOCX file.

    Args:
        source: Path to DOCX file

    Returns:
        str: Extracted text content
    """
    doc = docx.Document(source)
    text = "\n\n".join([paragraph.text for paragraph in doc.paragraphs])
    return text


def scrape_pptx(source: str) -> str:
    """
    Scrape text from a PPTX file.

    Args:
        source: Path to PPTX file

    Returns:
        str: Extracted text content
    """
    prs = pptx.Presentation(source)
    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n\n"

    return text


def scrape_image(source: str) -> str:
    """
    Scrape text from an image using OCR.

    Args:
        source: Path to image file

    Returns:
        str: Extracted text content
    """
    image = Image.open(source)
    text = pytesseract.image_to_string(image)
    image.close()
    return text


def scrape_audio(source: str, whisper_model_path: Optional[str] = None) -> str:
    """
    Transcribe audio file using MLX Whisper.

    Args:
        source: Path to audio file
        whisper_model_path: Path to MLX Whisper model

    Returns:
        str: Transcription
    """
    from mlxpipeline.core import DEFAULT_WHISPER_MODEL_PATH

    if not whisper_model_path:
        whisper_model_path = DEFAULT_WHISPER_MODEL_PATH

    try:
        # Load the Whisper model
        whisper = WhisperMLX(whisper_model_path)

        # Transcribe the audio
        result = whisper.transcribe(source)

        # Return the full text
        return result["text"]
    except Exception as e:
        print(f"Error transcribing audio: {str(e)}")
        return f"Error transcribing audio: {str(e)}"


def scrape_video(source: str, whisper_model_path: Optional[str] = None) -> str:
    """
    Extract audio from video and transcribe it using MLX Whisper.

    Args:
        source: Path to video file
        whisper_model_path: Path to MLX Whisper model

    Returns:
        str: Transcription
    """
    # Extract audio from video to a temporary file
    temp_file = tempfile.NamedTemporaryFile(suffix=".wav", delete=False)
    temp_file.close()

    try:
        # Use FFmpeg to extract audio
        subprocess.run(
            [
                "ffmpeg",
                "-i",
                source,
                "-ar",
                "16000",
                "-ac",
                "1",
                "-c:a",
                "pcm_s16le",
                temp_file.name,
            ],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
        )

        # Transcribe the audio
        transcription = scrape_audio(temp_file.name, whisper_model_path)

        return transcription
    except Exception as e:
        print(f"Error processing video: {str(e)}")
        return f"Error processing video: {str(e)}"
    finally:
        # Clean up
        if os.path.exists(temp_file.name):
            os.unlink(temp_file.name)


def scrape_markdown(source: str) -> str:
    """
    Extract text from a Markdown file, handling frontmatter.

    Args:
        source: Path to Markdown file

    Returns:
        str: Extracted text
    """
    try:
        # Parse frontmatter
        with open(source, "r", encoding="utf-8") as f:
            post = frontmatter.load(f)

        # Get content and metadata
        content = post.content
        metadata = post.metadata

        # Add metadata as a header if it exists
        result = ""
        if metadata:
            result += "---\n"
            for key, value in metadata.items():
                result += f"{key}: {value}\n"
            result += "---\n\n"

        result += content
        return result
    except Exception as e:
        # Fallback to simple reading
        print(f"Error parsing markdown with frontmatter: {str(e)}")
        with open(source, "r", encoding="utf-8") as f:
            return f.read()


def scrape_json(source: str) -> str:
    """
    Extract text from a JSON file.

    Args:
        source: Path to JSON file

    Returns:
        str: Formatted JSON text
    """
    with open(source, "r", encoding="utf-8") as f:
        data = json.load(f)

    # Format JSON as text
    return json.dumps(data, indent=2)


def scrape_csv(source: str) -> str:
    """
    Extract text from a CSV file.

    Args:
        source: Path to CSV file

    Returns:
        str: Formatted CSV text
    """
    result = []

    with open(source, "r", encoding="utf-8", newline="") as f:
        reader = csv.reader(f)
        for row in reader:
            result.append(", ".join(row))

    return "\n".join(result)


def scrape_webpage(url: str) -> str:
    """
    Scrape content from a webpage.

    Args:
        url: URL to scrape

    Returns:
        str: Extracted content
    """
    try:
        response = requests.get(
            url,
            headers={
                "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
                "AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.114 Safari/537.36"
            },
        )
        response.raise_for_status()

        # Try to detect encoding
        response.encoding = response.apparent_encoding

        # Parse HTML
        soup = BeautifulSoup(response.text, "html.parser")

        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.extract()

        # Extract text
        text = soup.get_text()

        # Clean up whitespace
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = "\n".join(chunk for chunk in chunks if chunk)

        return text
    except Exception as e:
        print(f"Error scraping webpage: {str(e)}")
        return f"Error scraping webpage: {str(e)}"


def scrape_youtube(url: str, whisper_model_path: Optional[str] = None) -> str:
    """
    Download and transcribe YouTube audio using MLX Whisper.

    Args:
        url: YouTube URL
        whisper_model_path: Path to MLX Whisper model

    Returns:
        str: Transcription
    """
    # Check if youtube-dl or yt-dlp is installed
    ytdl_command = None
    for cmd in ["yt-dlp", "youtube-dl"]:
        if shutil.which(cmd):
            ytdl_command = cmd
            break

    if not ytdl_command:
        return "Error: yt-dlp or youtube-dl not found. Please install one of them."

    # Create temporary directory
    temp_dir = tempfile.mkdtemp()
    audio_path = os.path.join(temp_dir, "audio.wav")

    try:
        # Download audio from YouTube
        subprocess.run(
            [ytdl_command, "--extract-audio", "--audio-format", "wav", "-o", audio_path, url],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
        )

        # Find the downloaded file (yt-dlp might add extensions)
        for file in os.listdir(temp_dir):
            if file.endswith(".wav"):
                audio_path = os.path.join(temp_dir, file)
                break

        # Transcribe the audio
        transcription = scrape_audio(audio_path, whisper_model_path)

        return transcription
    except Exception as e:
        print(f"Error processing YouTube video: {str(e)}")
        return f"Error processing YouTube video: {str(e)}"
    finally:
        # Clean up
        shutil.rmtree(temp_dir, ignore_errors=True)


def ai_extract_webpage_content(url: str, llm_model_path: Optional[str] = None) -> str:
    """
    Extract relevant content from a webpage using the local LLM.

    Args:
        url: URL to scrape
        llm_model_path: Path to MLX LLM model

    Returns:
        str: Extracted relevant content
    """
    from mlxpipeline.extract import extract_from_chunk

    # First, scrape the raw content
    raw_content = scrape_webpage(url)

    # Define a simple extraction schema for relevant content
    schema = """
    {
        "title": "string",
        "main_content": "string",
        "summary": "string"
    }
    """

    # Define an extraction prompt specific to webpages
    prompt = """
    I need to extract the main content from this webpage. Ignore navigation menus, 
    ads, footers, and other boilerplate. Focus on the main article, blog post, 
    or content that the page is primarily about.
    
    Here's the raw webpage content:
    {text}
    
    Extract the following information:
    {schema}
    
    Respond ONLY with the JSON object containing the extracted information.
    """

    # Extract relevant content
    result = extract_from_chunk(
        raw_content, schema=schema, extraction_prompt=prompt, llm_model_path=llm_model_path
    )

    # Parse the result and return the main content
    try:
        data = json.loads(result)
        if "main_content" in data and data["main_content"]:
            return data["main_content"]
        elif "summary" in data and data["summary"]:
            return data["summary"]
        else:
            return raw_content
    except:
        # Fallback to raw content if parsing fails
        return raw_content


def ai_extract_pdf_content(source: str, llm_model_path: Optional[str] = None) -> str:
    """
    Extract relevant content from a PDF using the local LLM.

    Args:
        source: Path to PDF file
        llm_model_path: Path to MLX LLM model

    Returns:
        str: Extracted relevant content
    """
    from mlxpipeline.chunker import chunk_text
    from mlxpipeline.extract import extract_from_chunk

    # First, scrape the raw content
    raw_content = scrape_pdf(source)

    # Chunk the content if it's too large
    if len(raw_content) > 10000:
        chunks = chunk_text(raw_content, chunk_size=10000, chunk_overlap=1000)
        # Use the first chunk for extraction
        content_to_extract = chunks[0]
    else:
        content_to_extract = raw_content

    # Define a simple extraction schema for relevant content
    schema = """
    {
        "title": "string",
        "main_content": "string",
        "summary": "string"
    }
    """

    # Define an extraction prompt specific to PDFs
    prompt = """
    I need to extract the main content from this PDF document. Ignore headers, 
    footers, page numbers, and other boilerplate. Focus on the main text, 
    key information, and content that the document is primarily about.
    
    Here's the raw PDF content:
    {text}
    
    Extract the following information:
    {schema}
    
    Respond ONLY with the JSON object containing the extracted information.
    """

    # Extract relevant content
    result = extract_from_chunk(
        content_to_extract, schema=schema, extraction_prompt=prompt, llm_model=llm_model
    )
